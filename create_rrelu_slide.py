#!/usr/bin/env python3
"""Create a single PPTX slide explaining RReLU with SFPI_Introduction background."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import os, tempfile, copy

SFPI_PPTX = "/localdev/vignjatijevic/tt-metal/SFPI_Introduction.pptx"
OUT = "/localdev/vignjatijevic/tt-metal/rrelu_slide.pptx"
TMP = tempfile.mkdtemp()
BG_IMG = os.path.join(TMP, "sfpi_bg.jpg")

# ── colours (dark text on light bg) ─────────────────────────────────
DARK = RGBColor(0x20, 0x20, 0x20)
ACCENT = RGBColor(0x5B, 0x21, 0xB6)  # deep purple (matches bg tint)
BODY = RGBColor(0x33, 0x33, 0x33)
SUBTEXT = RGBColor(0x66, 0x66, 0x66)


# ── extract background image from SFPI_Introduction ─────────────────
def extract_sfpi_bg():
    prs = Presentation(SFPI_PPTX)
    layout = prs.slide_layouts[9]  # "Content Page BG 1"
    blip = layout.background._element.find(".//" + qn("a:blip"))
    embed = blip.get(qn("r:embed"))
    image_part = layout.part.rels[embed].target_part
    with open(BG_IMG, "wb") as f:
        f.write(image_part.blob)


extract_sfpi_bg()


# ── helper: plot (transparent bg, dark-theme-friendly) ──────────────
def make_rrelu_plot(path):
    fig, ax = plt.subplots(figsize=(5.2, 3.6), dpi=180)
    fig.patch.set_alpha(0.0)
    ax.set_facecolor((1, 1, 1, 0.55))

    x = np.linspace(-4, 4, 500)
    lower, upper = 1.0 / 8, 1.0 / 3

    y_lo = np.where(x < 0, lower * x, x)
    y_hi = np.where(x < 0, upper * x, x)
    ax.fill_between(x, y_lo, y_hi, where=(x < 0), color="#7C3AED", alpha=0.18, label="slope band [l, u)")

    for a_val, c in [(lower, "#7C3AED"), (0.2, "#8B5CF6"), (upper, "#A78BFA")]:
        y = np.where(x < 0, a_val * x, x)
        ax.plot(x, y, color=c, linewidth=1.4, alpha=0.7)

    ax.plot(x[x >= 0], x[x >= 0], color="#DC2626", linewidth=2.5, label="x  (x ≥ 0)")

    y_relu = np.maximum(0, x)
    ax.plot(x, y_relu, color="#888888", linewidth=1.0, linestyle="--", alpha=0.5, label="ReLU")

    ax.axhline(0, color="#333", linewidth=0.4, alpha=0.3)
    ax.axvline(0, color="#333", linewidth=0.4, alpha=0.3)

    ax.set_xlabel("x", color="#202020", fontsize=11)
    ax.set_ylabel("RReLU(x)", color="#202020", fontsize=11)
    ax.set_title("RReLU  (lower=1/8, upper=1/3)", color="#202020", fontsize=12, fontweight="bold", pad=10)
    ax.tick_params(colors="#333", labelsize=9)
    for spine in ax.spines.values():
        spine.set_color("#999")

    leg = ax.legend(
        loc="upper left", fontsize=8, facecolor="white", edgecolor="#999", labelcolor="#333", framealpha=0.7
    )
    fig.tight_layout()
    fig.savefig(path, transparent=True, bbox_inches="tight")
    plt.close(fig)


def make_formula_img(path):
    fig, ax = plt.subplots(figsize=(5.0, 2.0), dpi=180)
    fig.patch.set_alpha(0.0)
    ax.set_facecolor((1, 1, 1, 0.0))
    ax.axis("off")

    lines = [
        ("RReLU(x)  =   x          if  x ≥ 0", 0.88, 14, "#202020", "bold"),
        ("                    a·x       if  x < 0", 0.72, 14, "#202020", "bold"),
        ("", 0.58, 10, "#202020", "normal"),
        ("Training:    a ~ Uniform(lower, upper)", 0.46, 12, "#5B21B6", "normal"),
        ("Eval:           a = (lower + upper) / 2", 0.26, 12, "#5B21B6", "normal"),
        ("Defaults:    lower = 1/8,  upper = 1/3", 0.08, 11, "#666666", "normal"),
    ]
    for text, y, size, color, weight in lines:
        ax.text(
            0.04,
            y,
            text,
            transform=ax.transAxes,
            fontsize=size,
            color=color,
            fontweight=weight,
            fontfamily="monospace",
            va="center",
        )

    ax.annotate(
        "",
        xy=(0.01, 0.72),
        xytext=(0.01, 0.88),
        arrowprops=dict(arrowstyle="-", color="#202020", lw=1.5),
        transform=ax.transAxes,
    )

    fig.savefig(path, transparent=True, bbox_inches="tight")
    plt.close(fig)


# ── generate images ──────────────────────────────────────────────────
plot_path = os.path.join(TMP, "rrelu_plot.png")
formula_path = os.path.join(TMP, "rrelu_formula.png")
make_rrelu_plot(plot_path)
make_formula_img(formula_path)

# ── build PPTX using SFPI as template ────────────────────────────────
# Use the SFPI presentation as a base so we inherit its slide master/layouts
prs = Presentation(SFPI_PPTX)

# Remove all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(qn("r:id"))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Add slide using "Content Page BG 1" layout (index 9) — inherits the bg
slide_layout = prs.slide_layouts[9]
slide = prs.slides.add_slide(slide_layout)

# Remove any placeholder shapes from the layout
for ph in list(slide.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

# ── title ────────────────────────────────────────────────────────────
tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "RReLU — Randomized Leaky Rectified Linear Unit"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK

# ── subtitle / source ───────────────────────────────────────────────
tx2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12), Inches(0.4))
tf2 = tx2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Source: pytorch.org/docs/stable/generated/torch.nn.RReLU.html"
p2.font.size = Pt(14)
p2.font.italic = True
p2.font.color.rgb = SUBTEXT

# ── bullet points (left column) ─────────────────────────────────────
bullets = [
    (
        "Stochastic negative slope",
        "For x < 0, slope a is sampled uniformly from [lower, upper)\n"
        "during training (defaults: lower=1/8, upper=1/3)",
    ),
    (
        "Deterministic at eval time",
        "At inference, a = (lower + upper) / 2 — no randomness,\n" "making the function fully reproducible",
    ),
    (
        "Regularization benefit",
        "Random slopes act as a built-in regularizer,\n" "reducing overfitting compared to fixed Leaky ReLU",
    ),
]

top = Inches(1.7)
for title_text, body_text in bullets:
    bx = slide.shapes.add_textbox(Inches(0.6), top, Inches(5.8), Inches(0.4))
    bf = bx.text_frame
    bf.word_wrap = True
    bp = bf.paragraphs[0]
    bp.text = "●  " + title_text
    bp.font.size = Pt(18)
    bp.font.bold = True
    bp.font.color.rgb = ACCENT

    bx2 = slide.shapes.add_textbox(Inches(1.0), top + Inches(0.40), Inches(5.4), Inches(0.9))
    bf2 = bx2.text_frame
    bf2.word_wrap = True
    bp2 = bf2.paragraphs[0]
    bp2.text = body_text
    bp2.font.size = Pt(14)
    bp2.font.color.rgb = BODY
    bp2.line_spacing = Pt(19)

    top += Inches(1.55)

# ── formula image (left-bottom) ─────────────────────────────────────
slide.shapes.add_picture(formula_path, Inches(0.4), Inches(5.45), width=Inches(5.6))

# ── plot image (right) ──────────────────────────────────────────────
slide.shapes.add_picture(plot_path, Inches(6.6), Inches(1.6), width=Inches(6.2))

# ── save ─────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved → {OUT}")
